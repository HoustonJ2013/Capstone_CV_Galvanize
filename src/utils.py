import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE,PP_ALIGN

def add_text(slide,left,top,text,fontsize=14):
    left_dis = left
    top_dis = top
    txtwid = Inches(2)
    txtheig = Inches(0.5)
    txbox = slide.shapes.add_textbox(left_dis, top_dis, txtwid, txtheig)
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(fontsize)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 255)
def strip_filename(filepath):
    return (filepath.split("/")[-1]).split(".")[0]
def max_side(w, h, max_side=4):
    if w >= h:
        width = Inches(max_side)
        height = width / w * h
    else:
        height = Inches(max_side)
        width = height / h * w
    return width, height
def add_4_pic_slide(prs, pics):
    blank_slide_layout = prs.slide_layouts[5]
    c_left_list = [0.2,5.5,0.2,5.5]
    c_top_list = [0.2,0.2,3.5,3.5]
    slide = prs.slides.add_slide(blank_slide_layout)
    for j in xrange(4):
        img_path = pics[j]
        imgname = strip_filename(img_path)
        im = Image.open(img_path)
        w, h = im.size
        c_left = Inches(c_left_list[j])
        c_top = Inches(c_top_list[j])
        width, height = max_side(w, h, max_side=3.2)
        res_text = imgname + str(w) + " x " + str(h)
        add_text(slide, c_left, c_top - Inches(0.2), res_text, fontsize=6)
        slide.shapes.add_picture(img_path, c_left, c_top, width=width, height=height)

def create_pptx(pptname = "test.pptx", image_list = None):
    prs = Presentation()
    n_pics = len(image_list)
    print(n_pics,n_pics/8)
    for i in xrange(n_pics/8):
        pic_names = [image_list[i * 8 + j * 2] for j in xrange(4)]
        add_4_pic_slide(prs,pic_names)
        pic_names = [image_list[i * 8 + j * 2 + 1] for j in xrange(4)]
        add_4_pic_slide(prs,pic_names)
    prs.save(pptname)

def jpg_image_to_array(image_path):
    with Image.open(image_path) as image:
        im_arr = np.fromstring(image.tobytes(), dtype=np.uint8)
        im_arr = im_arr.reshape((image.size[1], image.size[0], 3))
    return im_arr